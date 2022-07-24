from numpy.core.fromnumeric import shape

from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import MSO_AUTO_SIZE
from pprint import pprint
import urllib.parse as urlparse
from inspect import getmembers
from PIL import Image
import numpy as np
from requests.adapters import HTTPAdapter
from requests.packages.urllib3.util.retry import Retry
from types import FunctionType
import pandas as pd
from lxml import etree
import math
from pptx.dml.color import ColorFormat, RGBColor
import os
import six
import requests
import platform
import fnmatch
import copy
retry_strategy = Retry(
    total=3,
    status_forcelist=[429, 500, 502, 503, 504],
    method_whitelist=["HEAD", "GET", "OPTIONS"]
)
adapter = HTTPAdapter(max_retries=retry_strategy)
http = requests.Session()
http.mount("https://", adapter)
http.mount("http://", adapter)

# Spoof User-Agent for downloads
http = requests.Session()
http.headers.update({
    "User-Agent": "Mozilla/5.0 (X11; Ubuntu; Linux x86_64; rv:68.0) Gecko/20100101 Firefox/68.0"
})
def remove_shape(shape):
    ele = shape.element
    ele.getparent().remove(ele)
    
def copy_shape(slide, shape):
    el = shape.element
    newel = copy.deepcopy(el)
    slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

def clear_slide(slide):
    for shp in slide.shapes:
        remove_shape(shp)
def _get_blank_slide_layout(pres):
    layout_items_count = [len(layout.placeholders) for layout in pres.slide_layouts]
    min_items = min(layout_items_count)
    blank_layout_id = layout_items_count.index(min_items)
    return pres.slide_layouts[blank_layout_id]
def attributes(obj):
    disallowed_names = {
      name for name, value in getmembers(type(obj)) 
        if isinstance(value, FunctionType)}
    return {
      name: getattr(obj, name) for name in dir(obj) 
        if name[0] != '_' and name not in disallowed_names and hasattr(obj, name)}
def duplicate_slide(pres, index):
    # math.isnan(photo_url)
    
    """Duplicate the slide with the given index in pres.

    Adds slide to the end of the presentation"""
    source = pres.slides[index]

    blank_slide_layout = _get_blank_slide_layout(pres)
    dest = pres.slides.add_slide(blank_slide_layout)

    for shp in source.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        dest.shapes._spTree.insert_element_before(newel, 'p:extLst')

    for key, value in six.iteritems(source.part.rels):
        # Make sure we don't copy a notesSlide relation as that won't exist
        if not "notesSlide" in value.reltype:
            dest.part.rels.add_relationship(value.reltype, value._target, value.rId)
    left = top = Inches(0)
    img_path_back = "back.jpg"
    pic = dest.shapes.add_picture(img_path_back, left, top, width=prs.slide_width, height=prs.slide_height)

    # This moves it to the background
    dest.shapes._spTree.remove(pic._element)
    dest.shapes._spTree.insert(2, pic._element)
    
    return dest
def download_insert_image(slide, photo_url):
    if photo_url!='-1':
        if not os.path.exists("photos"):
            os.mkdir("photos")
        r = http.get(photo_url)
        if r.status_code == 200:
            with open("photos/"+os.path.basename(photo_url), 'wb') as f:
                f.write(r.content)
                f.close()
        img_path_src = "photos/"+os.path.basename(photo_url)
    else:
        img_path_src = "placeholder.png"
    index = 0
    for shape in slide.shapes:
        if shape.shape_type==13:
            # print("image index:",index)
            if index == 3:
                pic = slide.shapes.add_picture(img_path_src, shape.left, shape.top, shape.width, shape.height)
                slide.shapes._spTree.remove(shape._element)
                slide.shapes._spTree.remove(pic._element)
                slide.shapes._spTree.insert(3, pic._element)
            index = index + 1
def add_audio_button(slide, audio_url):
    if audio_url!='-1':
        if not os.path.exists("audios"):
            os.mkdir("audios")
        r = http.get(audio_url)
        if r.status_code == 200:
            with open("audios/"+os.path.basename(audio_url), 'wb') as f:
                f.write(r.content)
                f.close()
        audio_path = "audios/"+os.path.basename(audio_url)
        temp_audio_path = urlparse.urlparse(audio_url).path
        extension = os.path.splitext(temp_audio_path)[1]
        if extension  == ".mp3":
            mtype = "audio/mp3"
        elif extension == '.wav':
            mtype = 'audio/wav'
        elif extension == ".m4a":
            mtype = "audio/m4a"
        else:
            mtype = 'audio/wav'
        try:
            audio = slide.shapes.add_movie(audio_path, 
                0, 0, 30, 30, 
                mime_type = mtype, 
                poster_frame_image = "play_button.png")
            slide.shapes._spTree.remove(audio._element)
            slide.shapes._spTree.insert(9, audio._element)

            tree = audio._element.getparent().getparent().getnext().getnext()
            timing = [el for el in tree.iterdescendants() if etree.QName(el).localname == 'cond'][0]
            timing.set('delay', '0')

            print('Added Audio Button')


        except Exception as e:
            print('Creating Play Button failed!')
            print(e)


    
def copy_slide(src_slide, dst_slide):
    clear_slide(dst_slide)
    for shp in src_slide.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        dst_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

if __name__ == "__main__":
    template_file = "pvcc.pptx"
    output_file = "demo.pptx"
    data_file = "pvcc.xlsx"
    
    data = pd.read_excel(data_file)
    slide_count = len(data)
    prs = Presentation(template_file)
 
    for index, row in data.iterrows():
        # if index > 2:
        #     break
        temp_slide = duplicate_slide(prs, 0)
        photo = row['Photos']
        if pd.isnull(photo)==0:
            download_insert_image(temp_slide,photo)
        else:
            download_insert_image(temp_slide,"-1")
        audio_url = row['Audio']
        if pd.isnull(audio_url)==0:
            add_audio_button(temp_slide, audio_url)
        for shape in temp_slide.shapes:
            # print('%d %s' % (shape.shape_id, shape.name))
            if shape.shape_id == 4:
                shape.text = row['Diploma Name']
                shape.text_frame.paragraphs[0].font.size = Pt(45)
                shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            if shape.shape_id == 5:
                txt = ""
                if pd.isnull(row["Plan"])==0:
                    txt += row["Plan"] + "\n"
                if pd.isnull(row["Honors"])==0:
                    txt += row["Honors"]
                shape.text = txt
                shape.text_frame.paragraphs[0].font.size = Pt(37)
                shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(205, 205, 205)
            if shape.shape_id == 8:
                txt = ""
                if pd.isnull(row["Primary Degree"])==0:
                    txt = row["Primary Degree"]
                shape.text = txt
                shape.text_frame.paragraphs[0].font.size = Pt(33)
                shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    # cur_slide = prs.slides.add_slide(slide_layout)
    # copy_slide(temp_slide, cur_slide)
    #cur_slide.shapes.add_picture('graphic.jpg', Inches(6), Inches(1))  
    # for shape in cur_slide.shapes:
    #     if shape.shape_id == 56:
    #         shape.text = data['School'][0]
    #         #prs.save(output_file)
    #         print("id: %s, type: %s" % (shape.shape_id, shape.shape_type))
    #     if shape.shape_id == 57:
    #         #shape.text = "57"  font set -hometown
    #         #prs.save(output_file)
    #         print("id: %s, type: %s" % (shape.shape_id, shape.shape_type))
    #     if shape.shape_id == 58:
    #         #shape.text = "58" font set-quote
    #         #prs.save(output_file)
    #         print("id: %s, type: %s" % (shape.shape_id, shape.shape_type))
    #     if shape.shape_id == 59:
    #         shape.text = data['Hometown'][0]
    #         #prs.save(output_file)
    #         print("id: %s, type: %s" % (shape.shape_id, shape.shape_type))
    #     if shape.shape_id == 60:
    #         shape.text = data['Student statement'][0]
    #         #prs.save(output_file)
    #         print("id: %s, type: %s" % (shape.shape_id, shape.shape_type)) 
    #     if shape.shape_id == 63:
    #         shape.text = data['Diploma name'][0]
    #         #prs.save(output_file)
    #         print("id: %s, type: %s" % (shape.shape_id, shape.shape_type))
    #     if shape.shape_id == 64:
    #         #picture =shape
    #        # picture_ph = picture.insert_picture("Capture.png")
    #         #pic = prs.slides[0].shapes.add_picture("Capture.png", Inches(1), Inches[1])
    #        ## for shape in prs.slides[0].placeholders:
    #        #         if 'Picture' in shape.name:
    #        #             picture = prs.slides[0].placeholders[shape.placeholder_format.idx]
    #          #           picture.insert_picture("Capture.png")
    #         #placeholder = prs.slides[0].placeholders[1]
    #        # placeholder = placeholder.insert_picture("Capture.png")
    #         print("id: %s, type: %s" % (shape.shape_id, shape.shape_type))
    #     if shape.shape_id == 65:
    #         shape.text = data['Degree'][0]
    #         #prs.save(output_file)
    #         print("id: %s, type: %s" % (shape.shape_id, shape.shape_type))    
                                 
    prs.save(output_file)
    os.startfile(output_file)
   
   